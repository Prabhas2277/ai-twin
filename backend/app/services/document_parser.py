import os
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import google.generativeai as genai
from ..core.config import settings

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"Error parsing PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    text = ""
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"Error parsing DOCX {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")
    return text

def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"Error parsing TXT {file_path}: {e}")
    return ""

def extract_text_from_image_via_gemini(file_path: str) -> str:
    if not settings.GEMINI_API_KEY:
        return "[Error: GEMINI_API_KEY not configured. Cannot perform OCR on image.]"
    
    try:
        genai.configure(api_key=settings.GEMINI_API_KEY)
        # Using gemini-2.5-flash which is standard and supports image inputs
        model = genai.GenerativeModel("gemini-2.5-flash")
        
        # Read image bytes
        with open(file_path, "rb") as f:
            image_data = f.read()
            
        mime_type = "image/png"
        if file_path.lower().endswith((".jpg", ".jpeg")):
            mime_type = "image/jpeg"
            
        response = model.generate_content([
            {
                "mime_type": mime_type,
                "data": image_data
            },
            "Perform OCR on this image. Extract all text content verbatim. "
            "Maintain structural layout, tables, list items, and mathematical equations if present. "
            "Do not add introductions or commentaries, just output the extracted text."
        ])
        return response.text if response.text else ""
    except Exception as e:
        print(f"Error performing OCR via Gemini on {file_path}: {e}")
        return f"[Error: OCR failed: {str(e)}]"

def parse_document(file_path: str, file_type: str) -> str:
    """
    Dispatcher to parse documents depending on their extensions.
    """
    file_type = file_type.lower()
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return extract_text_from_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return extract_text_from_pptx(file_path)
    elif file_type == "txt":
        return extract_text_from_txt(file_path)
    elif file_type in ["png", "jpg", "jpeg", "webp"]:
        return extract_text_from_image_via_gemini(file_path)
    else:
        return ""
